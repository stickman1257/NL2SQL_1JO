from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pptx.oxml.ns as nsmap
from lxml import etree

# ── 색상 팔레트 (PPT 스타일) ──
NAVY        = RGBColor(0x1E, 0x3A, 0x5F)   # 헤더 배경
BLUE_MED    = RGBColor(0x2E, 0x75, 0xB6)   # 메인 파란색
BLUE_LIGHT  = RGBColor(0xBD, 0xD7, 0xEE)   # 연한 파란색 배경
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
GRAY_DARK   = RGBColor(0x40, 0x40, 0x40)
GRAY_MID    = RGBColor(0x80, 0x80, 0x80)
GRAY_LIGHT  = RGBColor(0xF2, 0xF2, 0xF2)
GRAY_LINE   = RGBColor(0xD0, 0xD0, 0xD0)

RED_STEP    = RGBColor(0xC0, 0x39, 0x2B)
ORANGE_STEP = RGBColor(0xE6, 0x7E, 0x22)
YELLOW_STEP = RGBColor(0xD4, 0xAC, 0x0D)
GREEN_STEP  = RGBColor(0x1E, 0x8B, 0x4C)
PURPLE_KB   = RGBColor(0x6C, 0x3A, 0xBD)
PURPLE_LITE = RGBColor(0xEB, 0xE5, 0xF8)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# 흰 배경
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE


# ════════════════════════════
# 헬퍼
# ════════════════════════════

def rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(1)):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(1)):
    s = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    s.adjustments[0] = 0.04
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        size=Pt(11), bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box

def multiline(slide, lines, l, t, w, h):
    """lines: [(text, size, bold, color, align)]"""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        text = item[0]
        size = item[1] if len(item) > 1 else Pt(10)
        bold = item[2] if len(item) > 2 else False
        color = item[3] if len(item) > 3 else BLACK
        align = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = size
        r.font.bold = bold
        r.font.color.rgb = color
    return box

def arrow_right(slide, x1, y, x2, color=GRAY_MID, w=Pt(1.5)):
    c = slide.shapes.add_connector(1, Inches(x1), Inches(y), Inches(x2), Inches(y))
    c.line.color.rgb = color
    c.line.width = w
    ln = c.line._ln
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    h = etree.SubElement(ln, f'{{{ns}}}headEnd')
    h.set('type', 'arrow')
    h.set('w', 'med')
    h.set('len', 'med')
    t = etree.SubElement(ln, f'{{{ns}}}tailEnd')
    t.set('type', 'none')
    return c

def arrow_down(slide, x, y1, y2, color=GRAY_MID, w=Pt(1.5)):
    c = slide.shapes.add_connector(1, Inches(x), Inches(y1), Inches(x), Inches(y2))
    c.line.color.rgb = color
    c.line.width = w
    ln = c.line._ln
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    h = etree.SubElement(ln, f'{{{ns}}}headEnd')
    h.set('type', 'arrow')
    h.set('w', 'med')
    h.set('len', 'med')
    t = etree.SubElement(ln, f'{{{ns}}}tailEnd')
    t.set('type', 'none')
    return c


# ════════════════════════════
# 1. 헤더 영역
# ════════════════════════════
# 네이비 상단 띠
rect(slide, 0, 0, 13.33, 1.2, fill=NAVY)

# 제목
txt(slide, "Schema Knowledge Base 자동 보강 파이프라인",
    0.4, 0.12, 10.0, 0.55,
    size=Pt(26), bold=True, color=WHITE)

txt(slide, "NL2SQL 실행 로그 → 실패 원인 분석 → LLM DB 탐색 → 컬럼 지식 누적 → SQL 생성 품질 향상",
    0.4, 0.68, 11.0, 0.38,
    size=Pt(11), color=BLUE_LIGHT)

# 오른쪽 서브 태그
txt(slide, "Qwen 2.5 3B  ·  ReAct Loop  ·  Daily Batch",
    10.1, 0.76, 3.1, 0.3,
    size=Pt(9), color=GRAY_MID, align=PP_ALIGN.RIGHT)


# ════════════════════════════
# 2. 왼쪽 메인 플로우 영역  (0.3 ~ 8.5)
# ════════════════════════════

# ── 섹션 타이틀 ──
txt(slide, "배치 파이프라인 (매일 1회 실행)",
    0.35, 1.32, 6.0, 0.32,
    size=Pt(10), bold=True, color=NAVY)
rect(slide, 0.35, 1.62, 7.9, 0.02, fill=NAVY)

# ── 입력 소스 ──
txt(slide, "입력", 0.35, 1.72, 1.0, 0.28, size=Pt(9), bold=True, color=GRAY_MID)

rrect(slide, 0.35, 2.02, 1.85, 0.72, fill=GRAY_LIGHT, line=GRAY_LINE)
txt(slide, "📋 JSONL 로그",    0.42, 2.08, 1.7, 0.28, size=Pt(9.5), bold=True, color=GRAY_DARK)
txt(slide, "NL2SQL 하루치\n실행 기록",  0.42, 2.33, 1.7, 0.38, size=Pt(8.5), color=GRAY_MID)

rrect(slide, 2.42, 2.02, 1.85, 0.72, fill=GRAY_LIGHT, line=GRAY_LINE)
txt(slide, "🗃️ SQLite DB",     2.49, 2.08, 1.7, 0.28, size=Pt(9.5), bold=True, color=GRAY_DARK)
txt(slide, "실제 데이터\n탐색 대상",   2.49, 2.33, 1.7, 0.38, size=Pt(8.5), color=GRAY_MID)

# 아래 방향 화살표
arrow_down(slide, 1.28, 2.74, 3.1, color=GRAY_MID)
arrow_down(slide, 3.34, 2.74, 3.1, color=GRAY_MID)


# ── STEP 박스 4개 ──
steps = [
    (RED_STEP,    "STEP 1",   "Hard Query Filter",
     "• difficulty = 'hard' 쿼리 추출\n• 실행 에러가 발생한 쿼리 포함\n• 보강 리소스를 어려운 쿼리에 집중"),

    (ORANGE_STEP, "STEP 2-3", "Column Selector",
     "• SQL 파싱 → table.column 참조 추출\n• alias 자동 해소 (o.status → orders.status)\n• 참조 빈도 집계, 가중 스코어 Top-N 선별"),

    (YELLOW_STEP, "STEP 4",   "DB Explorer Agent",
     "• Qwen 2.5 3B 로컬 LLM\n• ReAct 루프: 이유→실행→관찰 반복\n• SELECT 쿼리로 DB 직접 탐색"),

    (GREEN_STEP,  "STEP 5",   "Schema KB Updater",
     "• 탐색 결과 → KB JSON 저장\n• 버전 히스토리 누적 보존\n• DB 원본 메타데이터는 수정하지 않음"),
]

step_w = 1.82
gap    = 0.22
sx     = 0.35

for i, (color, num, title, desc) in enumerate(steps):
    bx = sx + i * (step_w + gap)

    # 스텝 박스 (상단 색상 띠 + 흰 본문)
    rrect(slide, bx, 3.1, step_w, 2.8, fill=WHITE, line=color, line_w=Pt(1.5))

    # 상단 색상 헤더
    hdr = rrect(slide, bx, 3.1, step_w, 0.62, fill=color)
    hdr.adjustments[0] = 0.04

    txt(slide, num,   bx, 3.14, step_w, 0.28, size=Pt(8.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, title, bx, 3.38, step_w, 0.32, size=Pt(11),  bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 설명 텍스트
    txt(slide, desc, bx + 0.1, 3.78, step_w - 0.18, 2.0,
        size=Pt(9), color=GRAY_DARK, wrap=True)

    # 화살표 (마지막 제외)
    if i < 3:
        ax = bx + step_w + 0.02
        arrow_right(slide, ax, 4.5, ax + gap - 0.02, color=GRAY_MID, w=Pt(2))

# ── 피드백 결과 박스 ──
rect(slide, 0.35, 6.05, 7.9, 0.02, fill=BLUE_MED)
rrect(slide, 0.35, 6.12, 7.9, 0.95, fill=BLUE_LIGHT, line=BLUE_MED)
txt(slide, "⬇  SQL 생성 시 자동 주입 (다음날 Agent가 참조)",
    0.5, 6.16, 6.0, 0.32, size=Pt(9.5), bold=True, color=NAVY)
txt(slide, '  - status  TEXT  [NOT_NULL]  |  Order lifecycle state',
    0.5, 6.44, 7.5, 0.26, size=Pt(9), color=GRAY_DARK)
txt(slide, '    enriched_note (v3, 2026-06-12): 값으로 \'cancelled\', \'refunded\', \'processing\' 등이 존재',
    0.5, 6.68, 7.5, 0.26, size=Pt(9), color=PURPLE_KB, italic=True)


# ════════════════════════════
# 구분선
# ════════════════════════════
rect(slide, 8.42, 1.32, 0.02, 5.75, fill=GRAY_LINE)


# ════════════════════════════
# 3. 오른쪽: KB 구조 설명
# ════════════════════════════
RX = 8.6
RW = 4.58

txt(slide, "Schema KB가 쌓이는 방식",
    RX, 1.32, RW, 0.32, size=Pt(10), bold=True, color=NAVY)
rect(slide, RX, 1.62, RW, 0.02, fill=NAVY)

# KB JSON 예시 박스
rrect(slide, RX, 1.72, RW, 2.45, fill=RGBColor(0xFA, 0xF8, 0xFF), line=PURPLE_KB)

multiline(slide, [
    ('  "orders.status" : {',           Pt(8.5), True,  PURPLE_KB),
    ('    "current_description":',       Pt(8.5), False, GRAY_DARK),
    ('      "\'cancelled\', \'refunded\',',  Pt(8.5), False, GRAY_DARK),
    ('       \'processing\' 등 존재",',  Pt(8.5), False, GRAY_DARK),
    ('    "history": [',                 Pt(8.5), False, GRAY_DARK),
    ('      { v1 · 06-03 · "TEXT형. 고유값 5개" },',   Pt(8),   False, GRAY_MID),
    ('      { v2 · 06-08 · "환불·취소 상태 포함" },',  Pt(8),   False, GRAY_MID),
    ('      { v3 · 06-12 · "현재 버전"  }  ← 오늘',   Pt(8),   True,  PURPLE_KB),
    ('    ]',                            Pt(8.5), False, GRAY_DARK),
    ('  }',                              Pt(8.5), True,  PURPLE_KB),
], RX + 0.12, 1.76, RW - 0.22, 2.35)

# 특징 5가지
feats = [
    ("🔄 버저닝",    "LLM이 잘못된 설명을 써도 이전 버전으로 롤백 가능"),
    ("📂 분리 저장", "원본 DB·CSV 메타데이터는 절대 수정하지 않음"),
    ("🎯 선별 기준", "참조 빈도 + 설명 빈약(30자↓) 컬럼에 가중치 ×1.5"),
    ("🔍 검색 강화", "KB 설명 토큰도 테이블 검색 점수에 반영"),
    ("💻 완전 로컬", "Qwen 2.5 3B 4bit — API 비용 0원, 1회 다운로드"),
]

fy = 4.3
for icon_title, desc in feats:
    rrect(slide, RX, fy, RW, 0.44, fill=GRAY_LIGHT, line=GRAY_LINE, line_w=Pt(0.5))
    txt(slide, icon_title, RX + 0.12, fy + 0.06, 1.6,  0.3, size=Pt(9), bold=True, color=NAVY)
    txt(slide, desc,        RX + 1.75, fy + 0.06, 3.3, 0.3, size=Pt(8.5), color=GRAY_DARK)
    fy += 0.5


# ── 저장 ──
out = "/Users/san/san_code/NL2SQL_1JO/assets/kb_enrichment_slide.pptx"
prs.save(out)
print(f"저장 완료: {out}")
